"""
HPTU Capstone Project - PowerPoint Generator
CSEE-611P | HPTU AI Assistant
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Constants ──────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x29, 0x80)   # HPTU navy blue
BLUE   = RGBColor(0x2D, 0x4E, 0xC2)   # accent blue
GOLD   = RGBColor(0xF7, 0x94, 0x1D)   # HPTU gold
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF0, 0xF2, 0xF8)   # light bg
DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # dark text
GRAY   = RGBColor(0x4A, 0x4A, 0x6A)   # secondary text

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # completely blank


def add_shape(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_bullet_box(slide, items, l, t, w, h,
                   size=16, color=DARK, bullet="▸ ", spacing=1.15):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after  = Pt(2)
        run = p.add_run()
        run.text = f"{bullet}{item}"
        run.font.size  = Pt(size)
        run.font.color.rgb = color
    return txb


def header_bar(slide, title, subtitle=None):
    """Navy header band across the top."""
    add_shape(slide, 0, 0, 13.33, 1.4, fill=NAVY)
    add_shape(slide, 0, 1.4, 13.33, 0.07, fill=GOLD)   # gold accent stripe
    add_text(slide, title, 0.4, 0.12, 12.5, 0.9,
             size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.85, 12.5, 0.5,
                 size=14, color=GOLD, align=PP_ALIGN.LEFT)


def footer_bar(slide, label="CSEE-611P | HPTU, Hamirpur"):
    add_shape(slide, 0, 7.1, 13.33, 0.4, fill=NAVY)
    add_text(slide, label, 0.3, 7.12, 9, 0.3,
             size=10, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, "© 2026 Kunal Chauhan", 10.5, 7.12, 2.5, 0.3,
             size=10, color=GOLD, align=PP_ALIGN.RIGHT)


def slide_bg(slide, color=LIGHT):
    add_shape(slide, 0, 0, 13.33, 7.5, fill=color)


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE SLIDE
# ══════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
slide_bg(s1, NAVY)

# Gold decorative bars
add_shape(s1, 0, 0, 0.25, 7.5, fill=GOLD)
add_shape(s1, 13.08, 0, 0.25, 7.5, fill=GOLD)
add_shape(s1, 0.25, 3.4, 12.83, 0.05, fill=GOLD)

# University name top
add_text(s1, "HIMACHAL PRADESH TECHNICAL UNIVERSITY", 0.5, 0.3, 12.3, 0.6,
         size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s1, "Hamirpur, Himachal Pradesh — 177001", 0.5, 0.85, 12.3, 0.4,
         size=12, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

# Separator line
add_shape(s1, 2.5, 1.3, 8.33, 0.04, fill=GOLD)

# Main title
add_text(s1, "HPTU AI ASSISTANT", 0.5, 1.5, 12.3, 1.0,
         size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s1, "An AI-Powered Smart Helpdesk & Career Consultant for HPTU Students",
         0.5, 2.4, 12.3, 0.7,
         size=17, color=GOLD, align=PP_ALIGN.CENTER, italic=True)

# Details box
add_shape(s1, 1.5, 3.55, 10.33, 2.85, fill=RGBColor(0x0F, 0x1B, 0x54))

add_text(s1, "Submitted By", 2.0, 3.65, 5, 0.4, size=12, color=GOLD, bold=True)
add_text(s1, "Kunal Chauhan", 2.0, 4.0, 5, 0.4, size=16, color=WHITE, bold=True)
add_text(s1, "Roll No: 23015103036", 2.0, 4.35, 5, 0.35, size=13, color=WHITE)
add_text(s1, "B.Tech — Computer Science & Engineering (6th Sem)", 2.0, 4.65, 5.5, 0.35, size=12, color=WHITE)

add_shape(s1, 7.3, 3.65, 0.04, 2.6, fill=GOLD)

add_text(s1, "Project Details", 7.6, 3.65, 5, 0.4, size=12, color=GOLD, bold=True)
add_text(s1, "Course Code: CSEE-611P", 7.6, 4.0, 5, 0.35, size=13, color=WHITE)
add_text(s1, "Capstone Project | Session 2025-26", 7.6, 4.35, 5, 0.35, size=13, color=WHITE)
add_text(s1, "Dept. of CSE | HPTU, Hamirpur", 7.6, 4.65, 5, 0.35, size=13, color=WHITE)

add_text(s1, "Guided By: Faculty Mentor, Dept. of CSE, HPTU", 0.5, 6.6, 12.3, 0.4,
         size=12, color=RGBColor(0xA0, 0xAE, 0xC0), align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — INTRODUCTION
# ══════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
slide_bg(s2)
header_bar(s2, "Introduction", "Background & Motivation")
footer_bar(s2)

add_text(s2, "What is HPTU AI Assistant?", 0.5, 1.65, 12, 0.45,
         size=18, bold=True, color=NAVY)
add_shape(s2, 0.5, 2.05, 1.5, 0.04, fill=GOLD)

bullets_intro = [
    "HPTU (Himachal Pradesh Technical University) serves thousands of students across HP with limited digital support infrastructure.",
    "Students frequently struggle to find: exam schedules, results, fee info, syllabus, previous year papers, and grievance contacts.",
    "Traditional university portals are static, hard to navigate, and do not provide personalized guidance.",
    "There is no single intelligent system that answers queries, resolves issues, and guides students in real time.",
    "This project delivers an AI-powered web assistant purpose-built for HPTU students and prospective applicants.",
]
add_bullet_box(s2, bullets_intro, 0.5, 2.2, 12.3, 4.6, size=15, color=DARK)

# side accent
add_shape(s2, 12.83, 1.5, 0.1, 5.5, fill=GOLD)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
slide_bg(s3)
header_bar(s3, "Problem Statement", "Challenges Faced by HPTU Students")
footer_bar(s3)

# Two columns
add_shape(s3, 0.5, 1.6, 5.9, 5.2, fill=WHITE)
add_shape(s3, 0.5, 1.6, 5.9, 5.2, fill=None, line=RGBColor(0xE2, 0xE6, 0xF0))
add_text(s3, "❌ Current Problems", 0.7, 1.7, 5.5, 0.5, size=15, bold=True, color=NAVY)
add_shape(s3, 0.7, 2.15, 1.0, 0.04, fill=GOLD)

probs = [
    "No single platform for all HPTU information",
    "Students miss important notices & deadlines",
    "Result/marksheet issues take weeks to resolve",
    "No career or placement guidance available",
    "PYQ papers scattered across multiple sites",
    "Exam date queries require manual website search",
    "No 24×7 support for student grievances",
]
add_bullet_box(s3, probs, 0.7, 2.3, 5.5, 4.2, size=13, color=DARK, bullet="✗  ")

add_shape(s3, 6.9, 1.6, 5.9, 5.2, fill=WHITE)
add_shape(s3, 6.9, 1.6, 5.9, 5.2, fill=None, line=RGBColor(0xE2, 0xE6, 0xF0))
add_text(s3, "✅ Proposed Solution", 7.1, 1.7, 5.5, 0.5, size=15, bold=True, color=NAVY)
add_shape(s3, 7.1, 2.15, 1.0, 0.04, fill=GOLD)

sols = [
    "Unified AI-powered assistant for all info",
    "Live notice scraping & real-time updates",
    "Instant result lookup by roll no. or name",
    "Career consultant with branch-wise guidance",
    "1000+ PYQ papers searchable via chatbot",
    "Live web search for exam dates & schedules",
    "24×7 helpdesk with contact routing",
]
add_bullet_box(s3, sols, 7.1, 2.3, 5.5, 4.2, size=13, color=DARK, bullet="✔  ")


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — METHODOLOGY: TOOLS & TECHNOLOGIES
# ══════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
slide_bg(s4)
header_bar(s4, "Methodology — Tech Stack", "Tools & Technologies Used")
footer_bar(s4)

tech_rows = [
    ("Backend",      "Python 3.12 + Flask 3.1 (Blueprints architecture)"),
    ("AI Engine",    "Cohere API — command-a-03-2025 / command-r7b-12-2024"),
    ("Database",     "MongoDB (pymongo 4.7) — notices, PYQs, chat logs, syllabus"),
    ("Web Scraping", "BeautifulSoup4 + Requests — himtu.ac.in & hptuonline.com"),
    ("PDF Processing","pypdf — text extraction from HPTU documents"),
    ("Scheduler",    "APScheduler 3.10 — auto-scraping every 30 minutes"),
    ("Frontend",     "HTML5 + Vanilla CSS3 + JavaScript (no frameworks)"),
    ("Deployment",   "Render (Gunicorn + Procfile) | Live: burgerbeast-projects.onrender.com"),
    ("Search",       "Google Custom Search API + DuckDuckGo fallback"),
    ("Speech",       "Web Speech API (en-IN locale) — voice input"),
]

col_w = 5.9
for i, (layer, detail) in enumerate(tech_rows):
    row, col = divmod(i, 2)
    x = 0.4 + col * (col_w + 0.6)
    y = 1.65 + row * 0.95
    add_shape(s4, x, y, col_w, 0.82, fill=WHITE)
    add_shape(s4, x, y, 0.08, 0.82, fill=GOLD)
    add_text(s4, layer, x + 0.2, y + 0.04, col_w - 0.3, 0.3, size=11, bold=True, color=NAVY)
    add_text(s4, detail, x + 0.2, y + 0.35, col_w - 0.3, 0.4, size=11, color=GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — METHODOLOGY: SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
slide_bg(s5)
header_bar(s5, "Methodology — System Architecture", "High-Level Design Overview")
footer_bar(s5)

# Architecture diagram (text-based boxes + arrows)
boxes = [
    (1.0,  1.65, 3.0, 0.8, "🌐 User Interface", "HTML5 + CSS3 + JS\n(Browser)"),
    (5.2,  1.65, 3.0, 0.8, "⚙️ Flask Backend", "Python 3.12\nBlueprint Routes"),
    (9.5,  1.65, 3.2, 0.8, "🧠 Cohere AI", "command-a-03-2025\nLLM Engine"),
    (1.0,  3.5,  3.0, 0.8, "🔍 Scraper", "Auto-scrapes himtu.ac.in\nEvery 30 mins"),
    (5.2,  3.5,  3.0, 0.8, "🗄️ MongoDB", "Notices, PYQ, Logs\nSyllabus, Fees"),
    (9.5,  3.5,  3.2, 0.8, "🕸️ Web Search", "Google CSE\n+ DuckDuckGo"),
    (3.2,  5.3,  3.0, 0.8, "📄 PDF Service", "pypdf — Text\nExtraction"),
    (7.0,  5.3,  3.0, 0.8, "📊 Result Scraper", "IndiaResults Portal\nLive Parsing"),
]

for (x, y, w, h, title, sub) in boxes:
    add_shape(s5, x, y, w, h, fill=NAVY)
    add_text(s5, title, x + 0.1, y + 0.05, w - 0.2, 0.38, size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s5, sub,   x + 0.1, y + 0.4,  w - 0.2, 0.38, size=10, color=WHITE, align=PP_ALIGN.CENTER)

# Arrow labels
add_text(s5, "→  HTTP Request", 4.05, 1.92, 1.1, 0.3, size=9, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s5, "→  AI Query",     8.25, 1.92, 1.2, 0.3, size=9, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s5, "↕  Read/Write",   4.05, 3.78, 1.1, 0.3, size=9, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s5, "→  Date Lookup",  8.25, 3.78, 1.2, 0.3, size=9, color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — METHODOLOGY: STEP-BY-STEP PROCESS
# ══════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
slide_bg(s6)
header_bar(s6, "Methodology — Process Flow", "Step-by-Step Working of the System")
footer_bar(s6)

steps = [
    ("01", "User Query Input",        "User types or speaks a query on the web interface. Voice input uses Web Speech API (en-IN)."),
    ("02", "Intent Detection",         "Backend classifies query: Result lookup, Date search, PYQ, Helpdesk issue, Career guidance, or General."),
    ("03", "Context Building",         "Smart context builder loads ONLY relevant data (notices, PYQ, fees, syllabus) from MongoDB cache."),
    ("04", "AI Processing",            "Cohere LLM processes the query + context and generates a concise, bullet-pointed Markdown response."),
    ("05", "Deterministic Overrides",  "For results/dates, deterministic scrapers override LLM to ensure factual accuracy (no hallucination)."),
    ("06", "Response Delivery",        "Formatted HTML response shown in chatbot UI with clickable links, emojis, and download buttons."),
]

for i, (num, title, desc) in enumerate(steps):
    y = 1.6 + i * 0.88
    add_shape(s6, 0.4, y, 0.7, 0.7, fill=GOLD)
    add_text(s6, num, 0.4, y + 0.1, 0.7, 0.5, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_shape(s6, 1.15, y, 11.7, 0.7, fill=WHITE)
    add_shape(s6, 1.15, y, 11.7, 0.7, fill=None, line=RGBColor(0xE2, 0xE6, 0xF0))
    add_text(s6, title, 1.35, y + 0.04, 3.5, 0.35, size=13, bold=True, color=NAVY)
    add_text(s6, desc,  4.9,  y + 0.04, 7.8,  0.62, size=12, color=GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — RESULTS: KEY FEATURES DELIVERED
# ══════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
slide_bg(s7)
header_bar(s7, "Results & Discussion", "Key Features Delivered")
footer_bar(s7)

features = [
    ("🤖", "AI Chatbot",        "Cohere-powered assistant answers 100+ query types about HPTU with Markdown-formatted, link-rich responses."),
    ("📊", "Result Lookup",     "Fetches live results from IndiaResults portal by roll number or student name with SGPA/CGPA/subject marks."),
    ("📄", "PYQ System",        "1000+ previous year papers scraped from hptuonline.com, searchable by course/branch/semester/subject."),
    ("📅", "Academic Calendar", "Live countdown timer + event cards showing official 2025-26 semester dates. AI answers calendar queries."),
    ("🎫", "Student Helpdesk",  "12+ issue types handled with exact steps, contact numbers, emails, and portal links for each problem type."),
    ("🎓", "Career Consultant", "Branch-wise career paths, GATE/CAT roadmaps, placement prep tips, semester-wise academic guidance."),
]

for i, (icon, title, desc) in enumerate(features):
    row, col = divmod(i, 2)
    x = 0.4 + col * 6.5
    y = 1.6 + row * 1.75
    add_shape(s7, x, y, 5.9, 1.55, fill=WHITE)
    add_shape(s7, x, y, 5.9, 1.55, fill=None, line=RGBColor(0xE2, 0xE6, 0xF0))
    add_shape(s7, x, y, 5.9, 0.07, fill=GOLD)
    add_text(s7, f"{icon} {title}", x + 0.2, y + 0.15, 5.5, 0.45, size=14, bold=True, color=NAVY)
    add_text(s7, desc, x + 0.2, y + 0.6, 5.5, 0.85, size=12, color=GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 8 — RESULTS: PERFORMANCE METRICS
# ══════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
slide_bg(s8)
header_bar(s8, "Results — Performance Metrics", "System Testing & Measurements")
footer_bar(s8)

metrics = [
    ("⚡", "AI Response Time",      "< 2.5 seconds",  "Average response time for general queries using command-r7b model"),
    ("📋", "Notices Scraped",       "50+ Notices",     "Auto-scraped from himtu.ac.in official notice board per cycle"),
    ("📄", "PYQ Papers Indexed",    "1000+ Papers",    "From hptuonline.com across 15+ courses and 8 B.Tech branches"),
    ("📚", "Documents Indexed",     "55+ Documents",   "PDFs extracted and made searchable for AI context"),
    ("🤖", "Query Types Handled",   "15+ Categories",  "Results, PYQ, fees, notices, syllabus, calendar, helpdesk, career"),
    ("🔄", "Auto-Refresh Cycle",    "Every 30 min",    "APScheduler-triggered scraping keeps data fresh automatically"),
    ("🎤", "Voice Input",           "100% Supported",  "Web Speech API integration with en-IN locale for Hindi-accented English"),
    ("📱", "Responsive Design",     "3 Breakpoints",   "Mobile, tablet, and desktop fully supported with dark/light mode"),
]

col_labels = ["Metric", "Value", "Details"]
for i, (icon, metric, value, detail) in enumerate(metrics):
    y = 1.62 + i * 0.67
    bg = RGBColor(0xF8, 0xF9, 0xFD) if i % 2 == 0 else WHITE
    add_shape(s8, 0.4, y, 12.5, 0.62, fill=bg)
    add_text(s8, f"{icon}  {metric}", 0.6,  y + 0.1, 3.8, 0.42, size=12, bold=True, color=NAVY)
    add_shape(s8, 4.5, y + 0.08, 1.8, 0.44, fill=GOLD)
    add_text(s8, value, 4.55, y + 0.1, 1.7, 0.42, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s8, detail, 6.45, y + 0.1, 6.3, 0.42, size=12, color=GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 9 — RESULTS: COMPARISON & DISCUSSION
# ══════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
slide_bg(s9)
header_bar(s9, "Results — Comparison & Discussion", "HPTU AI Assistant vs. Existing Solutions")
footer_bar(s9)

# Table header
headers = ["Feature", "HPTU Official Website", "Generic Chatbots", "HPTU AI Assistant ✅"]
widths  = [3.5, 2.5, 2.5, 3.0]
cols    = [0.3, 3.85, 6.35, 8.85]
y0 = 1.6

# Header row
for j, (hdr, w, x) in enumerate(zip(headers, widths, cols)):
    fill = NAVY if j == 0 else (BLUE if j < 3 else GOLD)
    add_shape(s9, x, y0, w, 0.55, fill=fill)
    add_text(s9, hdr, x + 0.1, y0 + 0.1, w - 0.15, 0.38,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Real-time AI Answers",    "❌ No",  "⚠️ Generic",  "✅ HPTU-specific"),
    ("Exam Result Lookup",       "⚠️ Portal only", "❌ No", "✅ In-chat lookup"),
    ("PYQ Papers",               "❌ No",  "❌ No",       "✅ 1000+ papers"),
    ("Student Helpdesk",         "❌ No",  "⚠️ Basic",    "✅ 12+ issue types"),
    ("Academic Calendar",        "⚠️ PDF only", "❌ No",  "✅ Live + AI-aware"),
    ("Career Guidance",          "❌ No",  "⚠️ Generic",  "✅ Branch-specific"),
    ("Voice Input",              "❌ No",  "⚠️ Some",     "✅ en-IN locale"),
    ("Auto-Updated Notices",     "✅ Yes", "❌ No",        "✅ Auto-scraped"),
]

for i, (feat, v1, v2, v3) in enumerate(rows):
    y = y0 + 0.55 + i * 0.6
    bg = RGBColor(0xF8, 0xF9, 0xFD) if i % 2 == 0 else WHITE
    for j, (val, w, x) in enumerate(zip([feat, v1, v2, v3], widths, cols)):
        add_shape(s9, x, y, w, 0.55, fill=bg)
        clr = NAVY if j == 0 else DARK
        add_text(s9, val, x + 0.1, y + 0.1, w - 0.15, 0.38,
                 size=11, bold=(j==0), color=clr, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 10 — CONCLUSION & FUTURE WORK
# ══════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
slide_bg(s10, NAVY)
add_shape(s10, 0, 0, 13.33, 7.5, fill=NAVY)

# Gold accents
add_shape(s10, 0, 0, 0.18, 7.5, fill=GOLD)
add_shape(s10, 13.15, 0, 0.18, 7.5, fill=GOLD)

header_bar(s10, "Conclusion & Future Work", "Summary & Roadmap")
add_shape(s10, 0, 1.4, 13.33, 0.07, fill=GOLD)

# Conclusion box
add_shape(s10, 0.5, 1.6, 6.0, 5.2, fill=RGBColor(0x0F, 0x1B, 0x54))
add_text(s10, "🏆 What We Achieved", 0.7, 1.7, 5.6, 0.5, size=15, bold=True, color=GOLD)
add_shape(s10, 0.7, 2.18, 1.2, 0.04, fill=GOLD)

conclusions = [
    "Built a full-stack AI web assistant specifically for HPTU",
    "Integrated Cohere LLM with smart, topic-aware context loading",
    "Live result lookup by roll number or student name",
    "1000+ PYQ papers searchable through natural language chat",
    "Official Academic Calendar with live countdown timer",
    "12+ student helpdesk issue types with exact resolution steps",
    "Branch-wise career consultant & placement guide",
    "Deployed live on Render with auto-scraping and MongoDB",
]
add_bullet_box(s10, conclusions, 0.7, 2.3, 5.6, 4.3, size=12, color=WHITE, bullet="✔ ")

# Future work box
add_shape(s10, 6.85, 1.6, 6.0, 5.2, fill=RGBColor(0x0F, 0x1B, 0x54))
add_text(s10, "🚀 Future Enhancements", 7.05, 1.7, 5.6, 0.5, size=15, bold=True, color=GOLD)
add_shape(s10, 7.05, 2.18, 1.4, 0.04, fill=GOLD)

future = [
    "Student login with personalised dashboard & history",
    "Hindi language support for regional accessibility",
    "Push notifications for new notices and deadlines",
    "Branch/college predictor based on HPCET/JEE rank",
    "Mobile app (React Native / Flutter)",
    "Integration with HPTU exam portal APIs (official)",
    "Formal student grievance ticketing system",
    "AI-powered study material recommendation engine",
]
add_bullet_box(s10, future, 7.05, 2.3, 5.6, 4.3, size=12, color=WHITE, bullet="▸ ")

# ── Save ──────────────────────────────────────────────────────
out_path = r"c:\Users\kunal\Downloads\capstone project\HPTU-AI-Assistant\HPTU_AI_Assistant_Capstone_PPT.pptx"
prs.save(out_path)
print(f"✅ Presentation saved: {out_path}")
